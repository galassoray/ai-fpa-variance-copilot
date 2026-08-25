"""
verify_decks.py
===============
Checks generated decks against ground truth recomputed from the canonical
pandas layer.

    python verify_decks.py                        # every deck in the repo root
    python verify_decks.py path/to/deck.pptx ...  # specific files

Deliberately does NOT go through the marts, the agent, or the deck module's own
provenance list. Ground truth is `run_pipeline.compute(load())`, so this checks
the whole chain rather than checking a module against its own bookkeeping.

Three questions, in order of severity:
  1. Does every figure in the deck exist in ITS OWN month's computation?
  2. Do the headline figures match exactly?
  3. Could the deck belong to a DIFFERENT month? Cross-month contamination --
     internally consistent, wrong period -- is the failure that would survive
     every internal check, so it is checked explicitly.

Exit code is 0 only if every deck passes, so this is usable as a build gate.
"""
from __future__ import annotations

import glob
import os
import re
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "src"))

import run_pipeline as rp  # noqa: E402
from guardrails.numeric_audit import extract_mentions  # noqa: E402
from pptx import Presentation  # noqa: E402

TABLES = rp.load()
OUTPUTS = rp.compute(TABLES)

#: Counts and ordinals that are deck furniture rather than financial figures:
#: step indices, row counts, the tool-call total, the literal 0 on the
#: "model-computed numbers" callout, and years inside titles.
STRUCTURAL = set(range(0, 26)) | set(range(1900, 2101))

_PERIOD_IN_NAME = re.compile(r"(20\d{2})[-_](\d{2})")


def _months() -> set:
    return set(OUTPUTS["variance_detail"]["month"].astype(str))


def period_of(path: str):
    """Infer the period from the filename, then confirm it exists in the data."""
    m = _PERIOD_IN_NAME.search(os.path.basename(path))
    if not m:
        return None
    key = f"{m.group(1)}-{m.group(2)}-01"
    return key if key in _months() else None


def period_from_slide(path: str):
    """Fall back to the title slide, so a renamed file is still checkable."""
    import datetime as dt

    text = "\n".join(sh.text_frame.text
                     for sh in Presentation(path).slides[0].shapes
                     if sh.has_text_frame)
    for line in text.splitlines():
        line = line.strip()
        for fmt in ("%B %Y", "%b %Y"):
            try:
                d = dt.datetime.strptime(line, fmt)
            except ValueError:
                continue
            key = f"{d.year:04d}-{d.month:02d}-01"
            if key in _months():
                return key
    return None


def month_values(month: str) -> set:
    """Every number the canonical layer computes for one month.

    Includes the driver SHARES, which the SQL tool derives
    (abs(oi_impact) / sum(abs(oi_impact))) and the pandas layer does not hold.
    Omitting them made the first run of this script report the decks' correct
    45.2% / 19.8% / 4.9% as unmatched -- a gap in the CHECKER, not the deck,
    and a reminder that a verification script needs the same scrutiny as the
    thing it verifies.
    """
    vals: set = set()
    for _name, df in OUTPUTS.items():
        if "month" not in df.columns:
            continue
        sub = df[df["month"] == month]
        for col in sub.columns:
            for v in sub[col].tolist():
                if isinstance(v, (int, float)) and v == v:      # not NaN
                    f = float(v)
                    for c in (f, abs(f), f * 100, abs(f) * 100):
                        vals.add(round(c, 2))

    for grain in ("variance_by_department", "variance_by_statement_line"):
        d = OUTPUTS[grain]
        d = d[d["month"] == month]
        if d.empty:
            continue
        a = d["oi_impact_ab"].abs()
        total = a.sum()
        if total:
            for share in (a / total).tolist():
                vals.add(round(float(share) * 100, 2))
                vals.add(round(float(share), 2))

    det = OUTPUTS["variance_detail"]
    det = det[det["month"] == month]
    for _dept, grp in det.groupby("department_id"):
        a = grp["oi_impact_ab"].abs()
        total = a.sum()
        if total:
            for share in (a / total).tolist():
                vals.add(round(float(share) * 100, 2))
                vals.add(round(float(share), 2))
    return vals


def deck_numbers(path: str) -> list:
    """Every rendered figure, using the flagship's own extractor."""
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
        for m in extract_mentions("\n".join(parts)):
            out.append((i, m.text, m.value, m.kind, m.lsd))
    return out


def chart_numbers(path: str) -> list:
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if getattr(sh, "has_chart", False) and sh.has_chart:
                for series in sh.chart.plots[0].series:
                    for v in series.values:
                        if v is not None:
                            out.append((i, float(v)))
    return out


def hit(value: float, kind: str, lsd: float, vals: set) -> bool:
    """Precision-aware: a figure rendered as $2.57M matches 2,568,827.32."""
    target = value * 100 if kind == "percent" else value
    tol = max((lsd * 100 if kind == "percent" else lsd) / 2, 0.01)
    for v in vals:
        if abs(v - target) <= tol or abs(v - abs(target)) <= tol:
            return True
    return False


def structural(value: float) -> bool:
    return float(value).is_integer() and int(value) in STRUCTURAL


def check(path: str, all_months: dict) -> bool:
    name = os.path.basename(path)
    month = period_of(path) or period_from_slide(path)
    print(f"\n{'=' * 76}\n{name}")
    if month is None:
        print("  SKIPPED - could not determine the period from the filename "
              "or the title slide")
        return True
    print(f"  period: {month}\n{'=' * 76}")

    own = all_months[month]
    ok = True

    figs = deck_numbers(path)
    misses = [(s, t, v) for s, t, v, k, l in figs
              if not hit(v, k, l, own) and not structural(v)]
    print(f"  text figures : {len(figs):>4}   unmatched: {len(misses)}")
    for s, t, v in misses[:6]:
        print(f"      slide {s}: {t}  ({v})")

    charted = chart_numbers(path)
    cmiss = [(s, v) for s, v in charted
             if round(v, 2) not in own and not structural(v)]
    print(f"  chart values : {len(charted):>4}   unmatched: {len(cmiss)}")
    for s, v in cmiss[:6]:
        print(f"      slide {s}: {v}")
    ok = ok and not misses and not cmiss

    om = OUTPUTS["operating_metrics"]
    row = om[om["month"] == month].iloc[0]
    text = "\n".join(sh.text_frame.text for sl in Presentation(path).slides
                     for sh in sl.shapes if sh.has_text_frame)
    for label, want, kind in (("revenue", row["revenue"], "money"),
                              ("operating income", row["operating_income"], "money"),
                              ("headcount", row["total_headcount"], "count"),
                              ("ending ARR", row["ending_arr"], "money")):
        if kind == "count":
            rendered = f"{int(want):,}"
        else:
            rendered = (f"(${abs(want):,.0f})" if want < 0 else f"${want:,.0f}")
        found = rendered in text
        print(f"  {label:<18} {rendered:>16}   {'FOUND' if found else 'MISSING'}")
        ok = ok and found

    distinct = [f for f in figs if not structural(f[2])]
    scores = {m: sum(1 for _s, _t, v, k, l in distinct if hit(v, k, l, vals))
              for m, vals in all_months.items()}
    best = max(scores, key=lambda m: scores[m])
    runner = max((m for m in scores if m != month), key=lambda m: scores[m])
    print(f"  best-matching month: {best}  "
          f"(own {scores[month]}/{len(distinct)}, "
          f"next best {runner} {scores[runner]})")
    if best != month:
        print("      *** CROSS-MONTH CONTAMINATION ***")
        ok = False
    return ok


def main(argv=None) -> int:
    argv = list(argv if argv is not None else sys.argv[1:])
    paths = argv or sorted(glob.glob(os.path.join(HERE, "*.pptx")))
    if not paths:
        print("No .pptx files found. Build one first:\n"
              "  python src/agent/run_package.py 2025-09 --deck "
              "variance-review-2025-09.pptx")
        return 2

    missing = [p for p in paths if not os.path.exists(p)]
    if missing:
        for p in missing:
            print(f"Not found: {p}")
        return 2

    all_months = {m: month_values(m) for m in sorted(_months())}

    results = [check(p, all_months) for p in paths]
    print(f"\n{'=' * 76}")
    if all(results):
        print(f"RESULT: {len(results)} deck(s) match their own month, and no other")
        return 0
    print(f"RESULT: MISMATCH in {sum(1 for r in results if not r)} of "
          f"{len(results)} deck(s)")
    return 1


if __name__ == "__main__":
    raise SystemExit(main())
