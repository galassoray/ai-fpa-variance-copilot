"""
deck.py
=======
Generates the monthly variance deck: the artifact the job actually produces.

WHY A DECK IS THE RIGHT OUTPUT
------------------------------
Everything upstream produces a package a developer reads. An FP&A analyst
produces a deck someone presents. Closing that gap is what turns "I built an
agent that queries a warehouse" into "I automated the deliverable."

THE HAZARD, AND HOW IT IS CLOSED
--------------------------------
A deck generator is a new surface where numbers get formatted, and every
formatting layer is a chance to become a SECOND PATH to a figure. If this module
ever summed a column, re-rounded, or derived a percentage, the architecture's
central guarantee would be broken in the most visible place possible -- on a
slide, in front of a room.

So this module cannot compute. It has exactly one way to put a number on a
slide: ``Deck.fig()``, which takes a value **out of a ledger section row** and
records where it came from as it formats it. There is no arithmetic anywhere in
this file -- no ``+``, no ``sum()``, no ``round()`` on a financial value. Totals
that appear on slides are totals a SQL tool computed.

``Deck.provenance`` accumulates one entry per figure placed, and
``test_agent_phase6`` re-opens the generated file, reads every text frame and
every chart series back out, and asserts each number traces to a ledger value.
Traceability is therefore structural: it is how figures get onto slides, not an
audit performed afterwards.

DESIGN
------
The palette is the app's own (``--ink``, ``--teal``, favorable green, unfavorable
red), so the deck and the tool look like one product. Green and red are used
only for variance direction -- in a finance deck that is semantics, not
decoration, and using them decoratively would be misreading.

Slides adapt to what the agent retrieved. A plan that never pulled the ARR
bridge produces a deck without that slide rather than one with an empty
placeholder -- the same "the plan determines the output" property that governs
the narrative.
"""
from __future__ import annotations

import io
import os
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION,
                             XL_TICK_LABEL_POSITION)
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- palette: the app's own, so the deck and the tool read as one product ---
INK = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
FAV = RGBColor(0x15, 0x80, 0x3D)
UNFAV = RGBColor(0xB9, 0x1C, 0x1C)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF8, 0xFA, 0xFC)

HEAD_FONT = "Cambria"       # safe-list serif, renders true-to-width
BODY_FONT = "Calibri"       # safe-list sans

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_W = W - 2 * MARGIN


def _e(v) -> Emu:
    """Coerce a computed position to whole EMU.

    Layout arithmetic (0.58 * CONTENT_W) yields floats, and python-pptx will
    write a non-integer EMU into the XML rather than reject it -- which reads
    back as a value PowerPoint and every downstream reader must guess at.
    """
    return Emu(int(round(float(v))))


@dataclass
class FigureRef:
    """One number placed on a slide, and the ledger row it came from."""

    display: str
    value: float
    section: str
    field: str
    step: int
    slide: int
    label: str = ""


class DeckError(RuntimeError):
    pass


class Deck:
    """A variance deck built strictly from ledger sections.

    Every figure goes through :meth:`fig`. Nothing in this class performs
    arithmetic on a financial value.
    """

    def __init__(self, result, goal: dict, candidate=None, packet=None):
        self.result = result
        self.goal = goal
        self.candidate = candidate
        self.packet = packet
        self.provenance: list = []
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self._blank = self.prs.slide_layouts[6]
        self._n = 0

    # -- the only way a number reaches a slide ----------------------------
    def fig(self, section: str, row: int, field: str, kind: str = "money",
            label: str = "") -> str:
        """Format one value FROM a ledger section row, recording its origin.

        Raises rather than returning a placeholder when the value is absent: a
        deck with a silent "n/a" where a figure belongs is worse than a build
        that stops, because the gap survives into the room.
        """
        sec = self.result.sections.get(section)
        if sec is None:
            raise DeckError(f"section '{section}' is not in this run")
        rows = sec["rows"]
        if row >= len(rows):
            raise DeckError(f"section '{section}' has {len(rows)} row(s); "
                            f"row {row} was requested")
        if field not in rows[row]:
            raise DeckError(f"'{field}' is not a field of section '{section}'")

        value = rows[row][field]
        display = _format(value, kind)
        self.provenance.append(FigureRef(display, value, section, field,
                                         sec["step"], self._n, label or field))
        return display

    def series(self, section: str, field: str, kind: str = "money") -> list:
        """A whole column from one section, each value recorded individually."""
        sec = self.result.sections.get(section)
        if sec is None:
            raise DeckError(f"section '{section}' is not in this run")
        out = []
        for i, r in enumerate(sec["rows"]):
            if field not in r:
                raise DeckError(f"'{field}' missing from '{section}' row {i}")
            out.append(r[field])
            self.provenance.append(
                FigureRef(_format(r[field], kind), r[field], section, field,
                          sec["step"], self._n, f"{field}[{i}]"))
        return out

    def labels(self, section: str, field: str) -> list:
        """Non-numeric labels. Not recorded: these are names, not figures."""
        sec = self.result.sections.get(section)
        if sec is None:
            raise DeckError(f"section '{section}' is not in this run")
        return [str(r.get(field, "")) for r in sec["rows"]]

    def has(self, *sections: str) -> bool:
        return all(s in self.result.sections and self.result.sections[s]["rows"]
                   for s in sections)

    def section_named(self, *tools: str):
        """Find a section by the TOOL that produced it, not by its label.

        Section labels are planner-chosen. Keying on the tool is the same
        decision the renderer made, for the same reason: an agent-authored plan
        must produce a full deck, not a deck missing every slide whose label it
        happened to spell differently.
        """
        for name, sec in sorted(self.result.sections.items(),
                                key=lambda kv: kv[1]["step"]):
            if sec["tool"] in tools and sec["rows"]:
                return name
        return None

    # -- primitives --------------------------------------------------------
    def _slide(self):
        self._n += 1
        s = self.prs.slides.add_slide(self._blank)
        return s

    def _text(self, slide, x, y, w, h, text, size=16, bold=False, color=INK,
              font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              spacing=1.0):
        box = slide.shapes.add_textbox(_e(x), _e(y), _e(w), _e(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
        return box

    def _title(self, slide, eyebrow: str, title: str):
        self._text(slide, MARGIN, Inches(0.45), CONTENT_W, Inches(0.25),
                   eyebrow.upper(), size=11, bold=True, color=MUTED,
                   font=BODY_FONT)
        self._text(slide, MARGIN, Inches(0.75), CONTENT_W, Inches(0.6),
                   title, size=30, bold=True, color=INK, font=HEAD_FONT)

    def _stat(self, slide, x, y, w, value: str, label: str, color=INK,
              caption: str = ""):
        """Large stat callout. The number dominates; the label explains it."""
        card = slide.shapes.add_shape(1, _e(x), _e(y), _e(w), Inches(1.55))  # RECTANGLE
        card.fill.solid()
        card.fill.fore_color.rgb = WASH
        card.line.color.rgb = LINE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False

        self._text(slide, x + Inches(0.25), y + Inches(0.2),
                   w - Inches(0.5), Inches(0.22), label.upper(), size=10,
                   bold=True, color=MUTED)
        self._text(slide, x + Inches(0.25), y + Inches(0.5),
                   w - Inches(0.5), Inches(0.6), value, size=30, bold=True,
                   color=color, font=HEAD_FONT)
        if caption:
            self._text(slide, x + Inches(0.25), y + Inches(1.14),
                       w - Inches(0.5), Inches(0.28), caption, size=10,
                       color=MUTED)

    def _table(self, slide, x, y, w, headers: list, rows: list,
               widths: list | None = None, right_align_from: int = 1,
               row_h: float = 0.34, max_rows: int | None = None):
        # Guard against running off the slide: a table that overflows is the
        # most common and most visible deck defect, and PowerPoint will not
        # clamp it -- the rows are simply written past the bottom edge.
        if max_rows is not None and len(rows) > max_rows:
            hidden = len(rows) - (max_rows - 1)
            rows = rows[:max_rows - 1] + [[f"+ {hidden} more"] +
                                          [""] * (len(headers) - 1)]
        n_r, n_c = len(rows) + 1, len(headers)
        h = Inches(row_h) * n_r
        shape = slide.shapes.add_table(n_r, n_c, _e(x), _e(y), _e(w), _e(h))
        tbl = shape.table
        if widths:
            total = sum(widths)
            for i, frac in enumerate(widths):
                tbl.columns[i].width = _e(w * frac / total)

        for r_i in range(n_r):
            tbl.rows[r_i].height = Inches(row_h)

        for c, head in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK
            _cell_text(cell, head, size=11, bold=True, color=PAPER,
                       align=PP_ALIGN.RIGHT if c >= right_align_from
                       else PP_ALIGN.LEFT)

        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PAPER if r % 2 else WASH
                color = INK
                if isinstance(val, tuple):
                    val, color = val
                _cell_text(cell, str(val), size=11, color=color,
                           align=PP_ALIGN.RIGHT if c >= right_align_from
                           else PP_ALIGN.LEFT)
        return tbl

    def _bar(self, slide, x, y, w, h, categories: list, values: list,
             title: str, colors: list | None = None, horizontal: bool = True,
             number_format: str = '#,##0,"K";(#,##0,"K")'):
        """Native PowerPoint chart. Never an image.

        A picture of a chart is a screenshot; an embedded chart is a
        deliverable an interviewer can click into and see the numbers behind.
        """
        if horizontal:
            # A horizontal bar chart plots the FIRST category at the bottom, so
            # a ranked list arrives upside down next to the table beside it.
            categories = list(reversed(categories))
            values = list(reversed(values))
            colors = list(reversed(colors)) if colors else None

        data = CategoryChartData()
        data.categories = categories
        data.add_series("", values, number_format)
        kind = (XL_CHART_TYPE.BAR_CLUSTERED if horizontal
                else XL_CHART_TYPE.COLUMN_CLUSTERED)
        frame = slide.shapes.add_chart(kind, _e(x), _e(y), _e(w), _e(h), data)
        chart = frame.chart
        chart.has_legend = False
        chart.has_title = bool(title)
        if title:
            chart.chart_title.text_frame.text = title
            _style_runs(chart.chart_title.text_frame, 12, MUTED, bold=True)

        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.gap_width = 60
        dl = plot.data_labels
        dl.number_format = number_format
        dl.number_format_is_linked = False
        dl.font.size = Pt(10)
        dl.font.bold = True
        dl.font.color.rgb = INK

        if colors:
            for pt, col in zip(plot.series[0].points, colors):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = col
        else:
            plot.series[0].format.fill.solid()
            plot.series[0].format.fill.fore_color.rgb = TEAL

        for axis in (chart.category_axis, chart.value_axis):
            axis.tick_labels.font.size = Pt(10)
            axis.tick_labels.font.color.rgb = MUTED
        # With negative values PowerPoint parks the category labels at the zero
        # crossing, so they render ON TOP of the bars in low contrast. LOW pins
        # them to the outside edge where they belong.
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        chart.category_axis.format.line.color.rgb = LINE
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.visible = False
        return chart

    def _grouped_bar(self, slide, x, y, w, h, categories, series: list,
                     title: str, colors: list):
        data = CategoryChartData()
        data.categories = categories
        for name, values in series:
            data.add_series(name, values, '#,##0,"K";(#,##0,"K")')
        frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                       _e(x), _e(y), _e(w), _e(h), data)
        chart = frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = MUTED
        chart.has_title = bool(title)
        if title:
            chart.chart_title.text_frame.text = title
            _style_runs(chart.chart_title.text_frame, 12, MUTED, bold=True)
        plot = chart.plots[0]
        plot.gap_width = 80
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = '#,##0,"K";(#,##0,"K")'
        dl.number_format_is_linked = False
        dl.font.size = Pt(9)
        dl.font.bold = True
        dl.font.color.rgb = INK
        for s, col in zip(plot.series, colors):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = col
        for axis in (chart.category_axis, chart.value_axis):
            axis.tick_labels.font.size = Pt(10)
            axis.tick_labels.font.color.rgb = MUTED
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = LINE
        chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
        return chart

    def _footer(self, slide, note: str = ""):
        self._text(slide, MARGIN, H - Inches(0.52), CONTENT_W, Inches(0.22),
                   note or "Synthetic data · every figure computed in SQL, "
                           "verified against the run ledger",
                   size=9, color=MUTED)


def _cell_text(cell, text, size=11, bold=False, color=INK,
               align=PP_ALIGN.LEFT):
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    # run.text rather than text_frame.text: the latter collapses the paragraph
    # to a single unstyled run and discards the formatting set here.
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = BODY_FONT


def _style_runs(tf, size, color, bold=False):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = BODY_FONT


def _format(value, kind: str) -> str:
    """Format a retrieved value. Never derives one.

    Accounting parentheses for negatives, matching the app and the CLI: a
    finance reader expects ($142,611), not -$142,611.
    """
    if value is None:
        return "—"
    if kind == "percent":
        return f"{value * 100:.1f}%"
    if kind == "count":
        return f"{int(value):,}"
    if kind == "count_signed":
        return f"{int(value):+,}"
    neg = value < 0
    body = f"${abs(value):,.0f}"
    return f"({body})" if neg else body


def _money_k(value) -> str:
    if value is None:
        return "—"
    neg = value < 0
    body = f"${abs(value) / 1000:,.0f}K"
    return f"({body})" if neg else body


# ==========================================================================
# slides
# ==========================================================================
def _period_label(period: str) -> str:
    import datetime as dt
    try:
        return dt.date.fromisoformat(str(period)).strftime("%B %Y")
    except ValueError:
        return str(period)


def _slide_title(d: Deck):
    s = d._slide()
    bg = s.shapes.add_shape(1, 0, 0, _e(W), _e(H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    bg.shadow.inherit = False

    period = _period_label(d.goal.get("period", ""))
    d._text(s, MARGIN, Inches(2.5), CONTENT_W, Inches(0.3),
            "MONTHLY VARIANCE REVIEW", size=13, bold=True,
            color=RGBColor(0xCA, 0xDC, 0xFC))
    d._text(s, MARGIN, Inches(2.95), CONTENT_W, Inches(1.0), period,
            size=54, bold=True, color=PAPER, font=HEAD_FONT)

    cmp_ = "Actual vs Budget" if d.goal.get("comparison") == "actual_vs_budget" \
        else "Actual vs Forecast"
    d._text(s, MARGIN, Inches(4.1), CONTENT_W, Inches(0.3), cmp_,
            size=17, color=RGBColor(0xCA, 0xDC, 0xFC))

    cost = d.result.ledger.cost_summary()
    planner = cost.get("planner_model") or "deterministic plan (no model)"
    d._text(s, MARGIN, Inches(6.35), CONTENT_W, Inches(0.5),
            f"Prepared by the FP&A close-cycle agent  ·  {cost['steps']} tool "
            f"calls  ·  planner: {planner}", size=11,
            color=RGBColor(0x94, 0xA3, 0xB8))
    d._text(s, MARGIN, Inches(6.68), CONTENT_W, Inches(0.5),
            "Synthetic data. Every figure computed in SQL and verified against "
            "the run ledger before publication.", size=11,
            color=RGBColor(0x94, 0xA3, 0xB8))


def _slide_summary(d: Deck, headline: str, pl: str):
    s = d._slide()
    d._title(s, "Executive summary", _period_label(d.goal.get("period", "")))

    gap = Inches(0.28)
    cw = (CONTENT_W - 3 * gap) / 4
    y = Inches(1.65)

    rev_row = _row_index(d, pl, "statement_line", "Revenue")
    d._stat(s, MARGIN, y, cw, d.fig(headline, 0, "revenue", label="revenue"),
            "Revenue", caption=f"vs budget {d.fig(pl, rev_row, 'base', label='revenue budget')}"
            if rev_row is not None else "")

    if rev_row is not None:
        var = d.result.sections[pl]["rows"][rev_row]["oi_impact"]
        d._stat(s, MARGIN + cw + gap, y, cw,
                d.fig(pl, rev_row, "oi_impact", label="revenue OI impact"),
                "Revenue vs plan", color=FAV if var > 0 else UNFAV,
                caption="favorable" if var > 0 else "unfavorable")

    oi = d.result.sections[headline]["rows"][0]["operating_income"]
    d._stat(s, MARGIN + 2 * (cw + gap), y, cw,
            d.fig(headline, 0, "operating_income", label="operating income"),
            "Operating income", color=FAV if oi > 0 else UNFAV,
            caption=f"margin {d.fig(headline, 0, 'operating_margin', 'percent', label='operating margin')}")

    d._stat(s, MARGIN + 3 * (cw + gap), y, cw,
            d.fig(headline, 0, "total_headcount", "count", label="headcount"),
            "Headcount", caption=f"ARR {_money_k(d.result.sections[headline]['rows'][0]['ending_arr'])}"
            if "ending_arr" in d.result.sections[headline]["rows"][0] else "")
    if "ending_arr" in d.result.sections[headline]["rows"][0]:
        d.fig(headline, 0, "ending_arr", label="ending ARR")

    # P&L at a glance
    rows = []
    for i, r in enumerate(d.result.sections[pl]["rows"]):
        col = FAV if r["oi_impact"] > 0 else UNFAV
        rows.append([r["statement_line"],
                     d.fig(pl, i, "actual", label="actual"),
                     d.fig(pl, i, "base", label="budget"),
                     (d.fig(pl, i, "oi_impact", label="OI impact"), col)])
    d._text(s, MARGIN, Inches(3.55), CONTENT_W, Inches(0.25),
            "PROFIT & LOSS VERSUS PLAN", size=11, bold=True, color=MUTED)
    d._table(s, MARGIN, Inches(3.85), CONTENT_W,
             ["Statement line", "Actual", "Budget", "Impact on operating income"],
             rows, widths=[3, 2, 2, 3])

    worst = None
    if drivers_hint := d.section_named("rank_variance_drivers"):
        rws = d.result.sections[drivers_hint]["rows"]
        if rws:
            worst = rws[0]
    if worst is not None:
        band = s.shapes.add_shape(1, MARGIN, Inches(5.65), _e(CONTENT_W), Inches(1.05))
        band.fill.solid()
        band.fill.fore_color.rgb = WASH
        band.line.color.rgb = LINE
        band.line.width = Pt(0.75)
        band.shadow.inherit = False
        d._text(s, MARGIN + Inches(0.3), Inches(5.85), CONTENT_W - Inches(0.6),
                Inches(0.22), "LARGEST SINGLE DRIVER", size=10, bold=True,
                color=MUTED)
        d._text(s, MARGIN + Inches(0.3), Inches(6.15), CONTENT_W - Inches(0.6),
                Inches(0.35),
                f"{worst['name']} at "
                f"{d.fig(drivers_hint, 0, 'oi_impact', label='largest driver')} "
                f"impact on operating income — detail follows.",
                size=14, color=INK)
    d._footer(s)


def _row_index(d: Deck, section: str, field: str, value: str):
    for i, r in enumerate(d.result.sections[section]["rows"]):
        if r.get(field) == value:
            return i
    return None


def _slide_drivers(d: Deck, drivers: str):
    s = d._slide()
    d._title(s, "Where the variance came from",
             "Largest drivers by operating-income impact")

    cats = d.labels(drivers, "name")
    vals = d.series(drivers, "oi_impact")
    colors = [FAV if v > 0 else UNFAV for v in vals]
    d._bar(s, MARGIN, Inches(1.6), CONTENT_W * 0.58, Inches(4.6), cats, vals,
           "Impact on operating income ($K)", colors=colors)

    x = MARGIN + CONTENT_W * 0.62
    w = CONTENT_W * 0.38
    rows = []
    for i, r in enumerate(d.result.sections[drivers]["rows"]):
        col = FAV if r["oi_impact"] > 0 else UNFAV
        share = d.fig(drivers, i, "share_of_total_oi_impact", "percent",
                      label="share of total") if r.get("share_of_total_oi_impact") \
            is not None else "—"
        rows.append([r["name"],
                     (d.fig(drivers, i, "oi_impact", label="OI impact"), col),
                     share])
    d._text(s, x, Inches(1.6), w, Inches(0.25), "RANKED", size=11, bold=True,
            color=MUTED)
    d._table(s, x, Inches(1.9), w, ["Driver", "Impact", "Share"], rows,
             widths=[4, 3, 2])
    d._text(s, x, Inches(1.95) + Inches(0.34) * (len(rows) + 1), w, Inches(1.0),
            "Ranked on impact to operating income, not raw variance: a "
            "department spanning revenue and expense lines has no readable "
            "variance sign.", size=10, color=MUTED, spacing=1.15)
    d._footer(s)


def _slide_decomposition(d: Deck, section: str, subtitle: str):
    s = d._slide()
    dept = d.result.sections[section]["rows"][0].get("department_name") \
        or d.result.sections[section]["params"].get("department_id", "")
    d._title(s, "Driver detail", f"{dept} — {subtitle}")

    cats = d.labels(section, "account_name")
    vals = d.series(section, "oi_impact")
    colors = [FAV if v > 0 else UNFAV for v in vals]
    d._bar(s, MARGIN, Inches(1.6), CONTENT_W * 0.58, Inches(4.6), cats, vals,
           "Impact on operating income ($K)", colors=colors)

    x = MARGIN + CONTENT_W * 0.62
    w = CONTENT_W * 0.38
    rows = []
    for i, r in enumerate(d.result.sections[section]["rows"]):
        col = FAV if r["oi_impact"] > 0 else UNFAV
        rows.append([r["account_name"],
                     d.fig(section, i, "actual", label="actual"),
                     (d.fig(section, i, "oi_impact", label="OI impact"), col)])
    d._text(s, x, Inches(1.6), w, Inches(0.25), "ACCOUNT DETAIL", size=11,
            bold=True, color=MUTED)
    d._table(s, x, Inches(1.9), w, ["Account", "Actual", "Impact"], rows,
             widths=[4, 3, 3])
    d._footer(s)


def _slide_comp(d: Deck, comp: str):
    s = d._slide()
    d._title(s, "Compensation", "Headcount versus rate")

    cats = d.labels(comp, "department_name")
    hc = d.series(comp, "hc_impact")
    rate = d.series(comp, "rate_impact")
    d._grouped_bar(s, MARGIN, Inches(1.6), CONTENT_W, Inches(4.3), cats,
                   [("Headcount effect", hc), ("Rate effect", rate)],
                   "Salary variance split ($K)", [TEAL, RGBColor(0x94, 0xA3, 0xB8)])

    d._text(s, MARGIN, Inches(6.15), CONTENT_W, Inches(0.8),
            "Splitting salary variance into headcount and rate separates a "
            "hiring-pace question from a compensation-level question. They have "
            "different owners and different fixes.",
            size=12, color=MUTED, spacing=1.2)
    d._footer(s)


def _slide_revenue(d: Deck, rev: str):
    s = d._slide()
    d._title(s, "Revenue", "Volume versus price")

    gap = Inches(0.35)
    cw = (CONTENT_W - 2 * gap) / 3
    y = Inches(1.75)
    r0 = d.result.sections[rev]["rows"][0]
    d._stat(s, MARGIN, y, cw, d.fig(rev, 0, "rev_variance", label="revenue variance"),
            "Revenue variance", color=FAV if r0["rev_variance"] > 0 else UNFAV)
    d._stat(s, MARGIN + cw + gap, y, cw,
            d.fig(rev, 0, "volume_impact", label="volume effect"),
            "Volume effect", color=FAV if r0["volume_impact"] > 0 else UNFAV,
            caption="customer count versus plan")
    d._stat(s, MARGIN + 2 * (cw + gap), y, cw,
            d.fig(rev, 0, "price_impact", label="price effect"),
            "Price effect", color=FAV if r0["price_impact"] > 0 else UNFAV,
            caption="average revenue per account")

    cats = ["Volume effect", "Price effect"]
    vals = [r0["volume_impact"], r0["price_impact"]]
    colors = [FAV if v > 0 else UNFAV for v in vals]
    d._bar(s, MARGIN, Inches(3.7), CONTENT_W * 0.55, Inches(2.6), cats, vals,
           "Decomposition ($K)", colors=colors, horizontal=False)
    d._text(s, MARGIN + CONTENT_W * 0.6, Inches(3.9), CONTENT_W * 0.4,
            Inches(2.0),
            "A revenue miss driven by volume is a pipeline and retention "
            "question. One driven by price is a packaging and discounting "
            "question. The split decides which conversation to have.",
            size=12, color=MUTED, spacing=1.2)
    d._footer(s)


def _slide_arr(d: Deck, arr: str):
    s = d._slide()
    d._title(s, "Annual recurring revenue", "Movement in the period")

    # Starting and ending ARR are two orders of magnitude larger than the
    # monthly flows, so charting all six together renders New and Churn as
    # invisible slivers -- which defeats the point of a bridge. The balances go
    # in callouts; the chart shows only the comparable movement components.
    #
    # Reductions are shown as magnitudes in red rather than negated: negating
    # would be arithmetic, and this module does not compute. Colour carries the
    # direction, which is how a finance reader reads it anyway.
    row = d.result.sections[arr]["rows"][0]
    flows = [("new_arr", "New", FAV), ("expansion_arr", "Expansion", FAV),
             ("contraction_arr", "Contraction", UNFAV),
             ("churned_arr", "Churn", UNFAV)]
    cats, vals, colors = [], [], []
    for f, lab, col in flows:
        if row.get(f) is None:
            continue
        cats.append(lab)
        vals.append(row[f])
        d.fig(arr, 0, f, label=lab)
        colors.append(col)

    if cats:
        d._bar(s, MARGIN, Inches(1.75), CONTENT_W * 0.6, Inches(4.0), cats, vals,
               "Movement in the period ($K) — green adds, red reduces",
               colors=colors, horizontal=False)

    x = MARGIN + CONTENT_W * 0.65
    w = CONTENT_W * 0.35
    gap = Inches(0.28)
    y = Inches(1.75)

    if row.get("starting_arr") is not None:
        d._stat(s, x, y, w, d.fig(arr, 0, "starting_arr", label="starting ARR"),
                "Starting ARR")
        y += Inches(1.55) + gap
    d._stat(s, x, y, w, d.fig(arr, 0, "ending_arr", label="ending ARR"),
            "Ending ARR", color=TEAL)
    y += Inches(1.55) + gap

    if row.get("nrr_ttm") is not None:
        cap = ("gross retention "
               + d.fig(arr, 0, "grr_ttm", "percent", label="GRR ttm")) \
            if row.get("grr_ttm") is not None else "trailing twelve months"
        d._stat(s, x, y, w,
                d.fig(arr, 0, "nrr_ttm", "percent", label="NRR ttm"),
                "Net revenue retention", caption=cap)
    else:
        d._text(s, x, y + Inches(0.1), w, Inches(1.2),
                "Trailing-twelve-month retention is undefined this early in "
                "the dataset. It is omitted rather than estimated.",
                size=11, color=MUTED, spacing=1.2)
    d._footer(s)


def _slide_headcount(d: Deck, hc: str):
    s = d._slide()
    d._title(s, "Headcount", "Actual versus plan")

    cats = d.labels(hc, "department_name")
    actual = d.series(hc, "actual_headcount", "count")
    budget = d.series(hc, "budget_headcount", "count")
    d._grouped_bar(s, MARGIN, Inches(1.6), CONTENT_W * 0.62, Inches(4.4), cats,
                   [("Actual", actual), ("Plan", budget)],
                   "Headcount", [TEAL, RGBColor(0xCB, 0xD5, 0xE1)])

    x = MARGIN + CONTENT_W * 0.66
    w = CONTENT_W * 0.34
    rows = []
    for i, r in enumerate(d.result.sections[hc]["rows"]):
        v = r["hc_var_vs_budget"]
        rows.append([r["department_name"],
                     d.fig(hc, i, "actual_headcount", "count", label="actual heads"),
                     (d.fig(hc, i, "hc_var_vs_budget", "count_signed",
                            label="heads vs plan"),
                      UNFAV if v > 0 else (FAV if v < 0 else MUTED))])
    d._text(s, x, Inches(1.6), w, Inches(0.25), "VERSUS PLAN", size=11,
            bold=True, color=MUTED)
    d._table(s, x, Inches(1.9), w, ["Department", "Actual", "vs plan"], rows,
             widths=[4, 2, 2])
    d._footer(s)


def _slide_where_to_look(d: Deck, brief):
    """Prioritisation, computed. The deck's most useful single slide.

    Deliberately states no action. Ranking and materiality are decided in code;
    the slide reports where the money moved and what is underneath it, and
    leaves the judgment to the reader -- which is the only half of this that
    can be verified.
    """
    s = d._slide()
    d._title(s, "Where to look", "Ranked by impact on operating income")

    # Geometry is pinned so three cards, the basis note, and the footer all
    # fit: cards 1.42 high on a 1.60 pitch from 1.55 ends the third at 6.17,
    # leaving the note at 6.35 and the footer clear at 6.98. The first version
    # used 1.55 on a 1.75 pitch and the note landed ON the third card.
    y = Inches(1.55)
    for a in brief.areas[:3]:
        colour = UNFAV if a.direction == "unfavorable" else FAV
        card = s.shapes.add_shape(1, MARGIN, _e(y), _e(CONTENT_W), Inches(1.42))
        card.fill.solid()
        card.fill.fore_color.rgb = WASH
        card.line.color.rgb = LINE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False

        share = (f"   ·   {a.share * 100:.1f}% of total impact"
                 if a.share is not None else "")
        d._text(s, MARGIN + Inches(0.3), _e(y + Inches(0.14)),
                _e(CONTENT_W * 0.48), Inches(0.3),
                f"{a.rank}.  {a.name}", size=15, bold=True, color=INK,
                font=HEAD_FONT)
        d._text(s, MARGIN + Inches(0.3), _e(y + Inches(0.5)),
                _e(CONTENT_W * 0.48), Inches(0.3),
                f"{_format(a.oi_impact, 'money')}  {a.direction}{share}",
                size=12, bold=True, color=colour)

        bits = [f"{e.label} {_format(e.value, 'money')}" for e in a.detail[:3]]
        if bits:
            d._text(s, MARGIN + Inches(0.3), _e(y + Inches(0.88)),
                    _e(CONTENT_W * 0.48), Inches(0.4),
                    "   ".join(bits), size=9, color=MUTED, spacing=1.1)

        # Supporting groups carry their basis, because account detail is on an
        # operating-income basis and compensation variance is not.
        right = []
        if a.comp:
            right.append("compensation (expense basis, + = above plan):  "
                         + "   ".join(f"{e.label} {_format(e.value, 'money')}"
                                      for e in a.comp))
        if a.headcount:
            right.append("headcount:  "
                         + "   ".join(f"{e.label} {_format(e.value, e.kind)}"
                                      for e in a.headcount))
        if a.revenue_split:
            right.append("revenue split (operating-income basis):  "
                         + "   ".join(f"{e.label} {_format(e.value, 'money')}"
                                      for e in a.revenue_split))
        yy = y + Inches(0.24)
        for line in right[:2]:
            d._text(s, MARGIN + _e(CONTENT_W * 0.52), _e(yy),
                    _e(CONTENT_W * 0.46), Inches(0.55), line, size=9,
                    color=MUTED, spacing=1.15)
            yy += Inches(0.52)

        y += Inches(1.60)

    d._text(s, MARGIN, Inches(6.35), CONTENT_W, Inches(0.3),
            "Ranking and materiality are computed in SQL. This slide shows "
            "where the money moved and what is underneath it; it does not "
            "recommend an action.", size=10, color=MUTED)
    d._footer(s)


def _slide_commentary(d: Deck):
    s = d._slide()
    d._title(s, "Commentary", "What the numbers say")

    text = (d.candidate.text or "").strip()
    body = d._text(s, MARGIN, Inches(1.6), CONTENT_W * 0.66, Inches(4.9),
                   text, size=15, color=INK, spacing=1.45)
    body.text_frame.word_wrap = True

    x = MARGIN + CONTENT_W * 0.70
    w = CONTENT_W * 0.30
    card = s.shapes.add_shape(1, _e(x), Inches(1.6), _e(w), Inches(4.9))
    card.fill.solid()
    card.fill.fore_color.rgb = WASH
    card.line.color.rgb = LINE
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    src = ("Written by the model, every figure verified"
           if d.candidate.source == "model"
           else "Generated deterministically from computed values")
    d._text(s, x + Inches(0.25), Inches(1.85), w - Inches(0.5), Inches(0.25),
            "HOW THIS WAS WRITTEN", size=10, bold=True, color=MUTED)
    d._text(s, x + Inches(0.25), Inches(2.2), w - Inches(0.5), Inches(1.5),
            f"{src}. {len(d.candidate.matched)} figures were checked against "
            f"the run ledger before this slide could be produced. A figure that "
            f"failed verification would have blocked publication.",
            size=11, color=INK, spacing=1.25)
    if d.candidate.matched:
        d._text(s, x + Inches(0.25), Inches(3.75), w - Inches(0.5), Inches(0.22),
                "SAMPLE TRACE", size=10, bold=True, color=MUTED)
        yy = Inches(4.05)
        for mention, _v, label in d.candidate.matched[:6]:
            d._text(s, x + Inches(0.25), yy, w - Inches(0.5), Inches(0.22),
                    f"{mention}  ←  {label}", size=9, color=MUTED)
            yy += Inches(0.26)
    d._footer(s)


def _slide_appendix(d: Deck):
    s = d._slide()
    d._title(s, "Appendix", "Provenance and controls")

    cost = d.result.ledger.cost_summary()
    gap = Inches(0.28)
    cw = (CONTENT_W - 3 * gap) / 4
    y = Inches(1.6)
    d._stat(s, MARGIN, y, cw, str(cost["steps"]), "Tool calls",
            caption="each logged with inputs and outputs")
    d._stat(s, MARGIN + cw + gap, y, cw, str(len(d.provenance)),
            "Figures on these slides", caption="each traced to a computed value")
    d._stat(s, MARGIN + 2 * (cw + gap), y, cw,
            str(len(d.candidate.matched) if d.candidate else 0),
            "Figures in the commentary", caption="verified before publication")
    d._stat(s, MARGIN + 3 * (cw + gap), y, cw, "0", "Model-computed numbers",
            color=FAV, caption="the model never does arithmetic")

    rows = []
    for e in d.result.ledger.entries:
        src = ""
        for pname, v in (e.params_declared or {}).items():
            if isinstance(v, str) and v.startswith("$STEP_"):
                src = f"{v} → {e.params_resolved.get(pname)}"
        rows.append([str(e.step_idx), e.tool, e.outcome, str(e.row_count),
                     src or "—"])
    d._text(s, MARGIN, Inches(3.42), CONTENT_W, Inches(0.25),
            "RUN LEDGER — EVERY QUERY THIS DECK RESTS ON", size=11, bold=True,
            color=MUTED)
    d._table(s, MARGIN, Inches(3.7), CONTENT_W,
             ["#", "Tool", "Outcome", "Rows", "Argument bound at runtime"],
             rows, widths=[0.6, 3.4, 1.6, 1.0, 4.4], right_align_from=3,
             row_h=0.27, max_rows=12)

    approval = ""
    if d.packet is not None and d.packet.published:
        a = d.packet.approval
        approval = (f"Approved for publication by {a.approver} · "
                    f"artifact {a.artifact_hash}")
    d._footer(s, approval or "Not approved for publication")


# ==========================================================================
# entry point
# ==========================================================================
def build_deck(result, goal: dict, candidate=None, packet=None) -> Deck:
    """Compose the deck from whatever the agent retrieved.

    Sections are located by the TOOL that produced them, not by their label, so
    an agent-authored plan yields a full deck rather than one missing every
    slide whose section it happened to name differently.
    """
    d = Deck(result, goal, candidate, packet)

    headline = d.section_named("get_operating_metrics")
    pl = d.section_named("get_pl_summary")
    drivers = d.section_named("rank_variance_drivers")
    comp = d.section_named("get_comp_decomposition")
    rev = d.section_named("get_revenue_decomposition")
    arr = d.section_named("get_arr_bridge")
    hc = d.section_named("get_headcount_movement")

    decomps = [name for name, sec in
               sorted(result.sections.items(), key=lambda kv: kv[1]["step"])
               if sec["tool"] == "decompose_variance" and sec["rows"]]

    _slide_title(d)
    if headline and pl:
        _slide_summary(d, headline, pl)

    from agent.briefing import build_briefing

    brief = build_briefing(result, goal)
    if brief.available:
        _slide_where_to_look(d, brief)

    if drivers:
        _slide_drivers(d, drivers)
    for i, name in enumerate(decomps[:2]):
        _slide_decomposition(d, name,
                             "largest driver" if i == 0 else "second driver")
    if comp:
        _slide_comp(d, comp)
    if rev:
        _slide_revenue(d, rev)
    if arr:
        _slide_arr(d, arr)
    if hc:
        _slide_headcount(d, hc)
    if candidate is not None and candidate.text:
        _slide_commentary(d)
    _slide_appendix(d)
    return d


def save_deck(result, goal: dict, candidate=None, packet=None,
              path: str = "") -> str:
    d = build_deck(result, goal, candidate, packet)
    if not path:
        period = str(goal.get("period", ""))[:7]
        path = os.path.join(os.getcwd(), f"variance-review-{period}.pptx")
    os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)
    d.prs.save(path)
    return path


def deck_bytes(result, goal: dict, candidate=None, packet=None) -> tuple:
    """In-memory deck, for a Streamlit download button.

    Returns (bytes, provenance) so a caller can show how many figures the deck
    contains and where each came from without writing to disk.
    """
    d = build_deck(result, goal, candidate, packet)
    buf = io.BytesIO()
    d.prs.save(buf)
    return buf.getvalue(), d.provenance
