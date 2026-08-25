"""
test_agent_phase6.py
====================
Phase-6 acceptance gate: the deck.

THE PROPERTY THIS SUITE EXISTS TO PROVE
----------------------------------------
**No number can appear on a slide unless a tool returned it.**

Not "we reviewed the deck and the numbers looked right" -- the file is re-opened
after generation, every text frame and every chart series is read back out, and
each number found is matched against a value in the run ledger. A figure that
matches nothing fails the build.

That matters because the deck is the most visible surface in the whole system.
Every previous layer earns its guarantee and then hands the numbers to a
formatter; if the formatter computes a total, re-rounds, or derives a
percentage, the guarantee dies on a slide, in a room, in front of the person
being persuaded by it.

The test reads the generated file rather than the generator's own bookkeeping,
because checking a module against its own records would only prove it is
self-consistent.
"""
import os
import re
import sys

import pytest

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "..", "src")
sys.path.insert(0, SRC)

pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402

import run_pipeline as rp  # noqa: E402
from agent import materialize as mz  # noqa: E402
from agent import tools as _tools  # noqa: E402,F401
from agent.deck import Deck, DeckError, build_deck, deck_bytes, save_deck  # noqa: E402
from agent.gates import PublicationPacket  # noqa: E402
from agent.narrate import narrate  # noqa: E402
from agent.orchestrator import Orchestrator  # noqa: E402
from agent.packages import build_goal, variance_package_plan  # noqa: E402
from agent.plan import Plan, Step  # noqa: E402
from guardrails import entity_audit as ea  # noqa: E402

PERIOD = "2025-09"


@pytest.fixture(scope="module")
def names():
    return ea.canonical_entity_names(rp.load())


@pytest.fixture(scope="module")
def con():
    mz.materialize(rp.compute(rp.load()), verbose=False)
    c = mz.connect_readonly()
    yield c
    c.close()


@pytest.fixture(scope="module")
def goal(con):
    return build_goal(con, PERIOD)


@pytest.fixture(scope="module")
def run(con, goal):
    return Orchestrator(con).run(variance_package_plan(goal), goal)


@pytest.fixture(scope="module")
def candidate(run, goal, names):
    return narrate(run, goal, None, names, mode="inject")


@pytest.fixture(scope="module")
def packet(run, candidate):
    p = PublicationPacket(run, candidate)
    p.approve("Test Reviewer", "acceptance run")
    return p


@pytest.fixture(scope="module")
def built(run, goal, candidate, packet):
    return build_deck(run, goal, candidate, packet)


@pytest.fixture(scope="module")
def pptx_path(tmp_path_factory, run, goal, candidate, packet):
    p = tmp_path_factory.mktemp("deck") / "variance.pptx"
    save_deck(run, goal, candidate, packet, str(p))
    return str(p)


# --------------------------------------------------------------------------
# helpers: read the FILE back, not the generator's bookkeeping
# --------------------------------------------------------------------------
#: The deck is checked with the SAME auditor that checks the prose.
#:
#: The first version of this test hand-rolled a number regex and immediately
#: got it wrong in the ways the flagship's extractor already solves: it split
#: "$28,809K" into "$28,809", read "$2.57M" as 2.57, and compared "45.2%"
#: against 45.164 without accounting for the precision it was displayed at.
#:
#: Reusing numeric_audit means the deck inherits the empirically derived
#: tolerance (MAX_REL_TOL = 0.005), the accounting-negative handling, and the
#: magnitude-suffix fix -- and it means a figure on a slide is held to exactly
#: the same standard as a figure in a sentence.
def _deck_text(prs) -> str:
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts)


def _slide_texts(prs) -> list:
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
        out.append((i, "\n".join(parts)))
    return out


def _chart_values(prs) -> list:
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False) and shape.has_chart:
                for series in shape.chart.plots[0].series:
                    for v in series.values:
                        if v is not None:
                            out.append((i, v))
    return out


def _ledger_values(result) -> set:
    """Every number any tool returned, plus its magnitude.

    Magnitudes are included because the deck renders negatives in accounting
    parentheses -- ($142,611) -- which parses back as -142611; and percentages
    are stored as ratios but rendered as percents.
    """
    vals = set()
    for rows in result.ledger.results.values():
        for row in rows:
            for v in row.values():
                if isinstance(v, (int, float)) and not isinstance(v, bool):
                    f = float(v)
                    vals.add(round(f, 2))
                    vals.add(round(abs(f), 2))
                    vals.add(round(f * 100, 2))       # ratio rendered as percent
                    vals.add(round(abs(f) * 100, 2))
    return vals


def _matches(value: float, ledger: set, tol: float = 0.02) -> bool:
    r = round(value, 2)
    if r in ledger:
        return True
    # Rendered figures are rounded for display ($2,568,827 from 2568827.32;
    # 81.0% from 0.8101). Accept a value that rounds to a ledger value at the
    # precision it was displayed with.
    for cand in (r, abs(r)):
        for lv in ledger:
            if abs(lv - cand) <= max(tol, abs(lv) * 5e-4):
                return True
    return False


#: Counts and ordinals that are structurally part of the deck rather than
#: financial figures: step indices, row counts, the tool-call total, the
#: literal 0 on the "model-computed numbers" callout, and years inside labels.
def _structural(value: float, result) -> bool:
    if float(value).is_integer():
        v = int(value)
        if 0 <= v <= max(len(result.ledger.entries), 24) + 1:
            return True
        if 1900 <= v <= 2100:                     # a year in a title
            return True
        if v == len(result.sections):
            return True
    return False


# --------------------------------------------------------------------------
# 1. the property
# --------------------------------------------------------------------------
def test_every_number_on_every_slide_traces_to_the_ledger(pptx_path, run, goal):
    """The whole point of the deck module's design.

    Reads the generated FILE and audits it with the same numeric auditor used
    on the commentary: checking a module against its own provenance list would
    prove only that it is self-consistent.
    """
    from agent.facts import fact_pack_from_ledger
    from guardrails import numeric_audit as na

    pack = fact_pack_from_ledger(run, goal)
    prs = Presentation(pptx_path)

    failures = []
    for slide_no, text in _slide_texts(prs):
        if not text.strip():
            continue
        res = na.audit(text, pack)
        for v in res.violations:
            # Structural counts -- step indices, row counts, the tool-call
            # total, the literal 0 on "model-computed numbers" -- are part of
            # the deck's furniture, not financial figures.
            if _structural(v.value, run):
                continue
            failures.append((slide_no, v.mention, v.value, v.reason))

    assert not failures, (
        "figures on slides that no tool returned: "
        + "; ".join(f"slide {s}: {m} ({r})" for s, m, _v, r in failures[:8])
    )


def test_every_charted_value_traces_to_the_ledger(pptx_path, run):
    """Charts are the easiest place for a derived number to hide, because the
    value never appears as text a reader could check."""
    prs = Presentation(pptx_path)
    ledger = set()
    for rows in run.ledger.results.values():
        for row in rows:
            for v in row.values():
                if isinstance(v, (int, float)) and not isinstance(v, bool):
                    ledger.add(round(float(v), 2))
                    ledger.add(round(abs(float(v)), 2))

    unexplained = [(s, v) for s, v in _chart_values(prs)
                   if round(float(v), 2) not in ledger
                   and not _structural(v, run)]
    assert not unexplained, (
        "charted values with no ledger source: "
        + "; ".join(f"slide {s}: {v}" for s, v in unexplained[:10])
    )


def test_the_generator_contains_no_arithmetic_on_figures():
    """Enforced by reading the source, so it survives future edits.

    A formatter that sums a column or derives a percentage becomes a second
    path to a number, which is exactly what every layer below this one exists
    to prevent.
    """
    src = open(os.path.join(SRC, "agent", "deck.py"), encoding="utf-8").read()

    # Layout arithmetic above the slide section is fine and unavoidable --
    # column widths, card positions. What must never happen is arithmetic on a
    # retrieved FIGURE, and that only occurs where slides are composed.
    marker = "# slides"
    assert marker in src
    slides_src = src[src.index(marker):]
    body = "\n".join(line for line in slides_src.splitlines()
                     if not line.strip().startswith("#"))
    # Word-boundary matched: a bare substring search reports "round(" inside
    # "background(", which is how this test first failed.
    banned = r"\b(sum|mean|round|abs|min|max)\s*\(|\bstatistics\.|\bnp\.|\bmath\."
    hits = [line.strip() for line in body.splitlines()
            if re.search(banned, line)]
    assert not hits, (
        "the slide layer performs arithmetic; every figure must come from "
        f"Deck.fig() or Deck.series(). Offending lines: {hits[:3]}"
    )


def test_provenance_records_one_entry_per_figure(built):
    assert built.provenance
    for ref in built.provenance:
        assert ref.section in built.result.sections
        assert ref.step > 0 and ref.slide > 0
        assert ref.display


def test_fig_refuses_a_value_that_was_not_retrieved(run, goal):
    d = Deck(run, goal)
    with pytest.raises(DeckError, match="not in this run"):
        d.fig("no_such_section", 0, "actual")
    section = sorted(run.sections)[0]
    with pytest.raises(DeckError, match="row 99 was requested"):
        d.fig(section, 99, "actual")
    with pytest.raises(DeckError, match="not a field"):
        d.fig(section, 0, "invented_field")


# --------------------------------------------------------------------------
# 2. the deck is well-formed
# --------------------------------------------------------------------------
def test_the_deck_opens_and_has_the_expected_shape(pptx_path):
    prs = Presentation(pptx_path)
    assert len(prs.slides) >= 9
    assert prs.slide_width > prs.slide_height     # 16:9


def test_charts_are_native_not_images(pptx_path):
    """A picture of a chart is a screenshot; an embedded chart is a deliverable
    an interviewer can click into and see the numbers behind."""
    prs = Presentation(pptx_path)
    charts = sum(1 for s in prs.slides for sh in s.shapes
                 if getattr(sh, "has_chart", False) and sh.has_chart)
    pictures = sum(1 for s in prs.slides for sh in s.shapes
                   if sh.shape_type is not None and "PICTURE" in str(sh.shape_type))
    assert charts >= 5, f"expected native charts, found {charts}"
    assert pictures == 0, "charts must not be rasterised into images"


def test_no_text_overflows_its_shape(pptx_path):
    """Text spilling past its box is the most common and most visible deck
    defect, and PowerPoint does not clamp it."""
    prs = Presentation(pptx_path)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            assert shape.left >= 0 and shape.top >= 0, f"slide {i}: negative offset"
            assert shape.left + (shape.width or 0) <= slide_w + 1000, (
                f"slide {i}: shape runs off the right edge")
            assert shape.top + (shape.height or 0) <= slide_h + 1000, (
                f"slide {i}: shape runs off the bottom edge")


def test_the_appendix_reports_the_provenance_counts(pptx_path, built):
    prs = Presentation(pptx_path)
    text = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                    if sh.has_text_frame)
    assert "provenance and controls" in text.lower()
    assert "the model never does arithmetic" in text
    assert str(len(built.result.ledger.entries)) in text


def test_the_deck_states_the_data_is_synthetic(pptx_path):
    """A clearance line, not a caption. It must appear without being sought."""
    prs = Presentation(pptx_path)
    text = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                    if sh.has_text_frame).lower()
    assert "synthetic" in text


def test_approval_appears_in_the_appendix(pptx_path, packet):
    prs = Presentation(pptx_path)
    text = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                    if sh.has_text_frame)
    assert "Test Reviewer" in text and packet.approval.artifact_hash in text


# --------------------------------------------------------------------------
# 3. the deck follows the plan
# --------------------------------------------------------------------------
def _thin_run(con, goal, *steps):
    plan = Plan(goal="g",
                steps=[Step(i, t, p, purpose=f"s{i}")
                       for i, (t, p) in enumerate(steps, start=1)],
                promises=[f"s{i}" for i in range(1, len(steps) + 1)])
    return Orchestrator(con).run(plan, goal)


def test_a_thin_plan_produces_a_shorter_deck_not_an_empty_slide(con, goal, names,
                                                                tmp_path):
    """Same property as the narrative: the plan determines the output. A plan
    that never pulled the ARR bridge yields a deck without that slide rather
    than one with a placeholder."""
    thin = _thin_run(con, goal,
                     ("get_operating_metrics", {"period": "$GOAL.period"}),
                     ("get_pl_summary", {"period": "$GOAL.period"}))
    cand = narrate(thin, goal, None, names, mode="inject")
    path = str(tmp_path / "thin.pptx")
    save_deck(thin, goal, cand, None, path)

    prs = Presentation(path)
    text = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                    if sh.has_text_frame)
    low = text.lower()
    assert "annual recurring revenue" not in low
    assert "headcount versus rate" not in low
    assert "executive summary" in low


def test_an_agent_authored_plan_still_produces_a_full_deck(con, goal, names,
                                                           tmp_path):
    """Sections are located by TOOL, not by the planner's own labels -- the
    same decision the renderer made, for the same reason."""
    agent = _thin_run(
        con, goal,
        ("get_pl_summary", {"period": "$GOAL.period"}),
        ("rank_variance_drivers",
         {"period": "$GOAL.period", "dimension": "department", "top_n": 5}),
        ("decompose_variance",
         {"period": "$GOAL.period", "department_id": "$STEP_2.rows[0].member",
          "top_n": 5}),
        ("get_operating_metrics", {"period": "$GOAL.period"}),
    )
    cand = narrate(agent, goal, None, names, mode="inject")
    path = str(tmp_path / "agent.pptx")
    save_deck(agent, goal, cand, None, path)

    prs = Presentation(path)
    text = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes
                    if sh.has_text_frame)
    low = text.lower()
    for expected in ("executive summary", "where the variance came from",
                     "driver detail", "commentary", "appendix"):
        assert expected in low, f"missing '{expected}' for an agent-authored plan"


@pytest.mark.parametrize("period", ["2024-06", "2025-03", "2025-12"])
def test_the_deck_builds_for_every_period(con, names, period, tmp_path):
    """2024-06 has no trailing-twelve-month retention -- the edge case that
    previously raised in the narrative layer."""
    g = build_goal(con, period)
    res = Orchestrator(con).run(variance_package_plan(g), g)
    cand = narrate(res, g, None, names, mode="inject")
    path = str(tmp_path / f"{period}.pptx")
    save_deck(res, g, cand, None, path)
    assert Presentation(path).slides


def test_deck_bytes_returns_a_valid_file_and_its_provenance(run, goal, candidate):
    data, prov = deck_bytes(run, goal, candidate, None)
    assert data[:2] == b"PK" and len(data) > 20_000
    assert prov and all(p.section in run.sections for p in prov)


# --------------------------------------------------------------------------
# 4. CLI and app integration
# --------------------------------------------------------------------------
def test_run_package_deck_flag(tmp_path, capsys, con):
    from agent import run_package as cli

    out = str(tmp_path / "cli.pptx")
    assert cli.main([PERIOD, "--deck", out]) == 0
    assert "Deck written" in capsys.readouterr().out
    assert Presentation(out).slides


def test_run_agent_baseline_deck_flag(tmp_path, capsys, con, monkeypatch):
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    from agent import run_agent

    out = str(tmp_path / "agent.pptx")
    assert run_agent.main(["--period", PERIOD, "--baseline-only",
                           "--deck", out]) == 0
    assert "each traced to a computed value" in capsys.readouterr().out
    assert Presentation(out).slides


def test_the_app_page_renders_the_package_visually_and_offers_the_deck(monkeypatch):
    """The page used to print a monospace block, which made a finished tool
    read like a terminal. It now renders metrics and tables, and offers the
    deck the job actually produces."""
    pytest.importorskip("streamlit")
    from streamlit.testing.v1 import AppTest

    for env in ("OPENAI_API_KEY", "ANTHROPIC_API_KEY"):
        monkeypatch.delenv(env, raising=False)
    sys.path.insert(0, os.path.join(HERE, "..", "eval"))

    at = AppTest.from_file(os.path.join(HERE, "..", "app", "app.py"),
                           default_timeout=240)
    at.run()
    at = at.sidebar.radio[0].set_value("Close-cycle agent").run()
    assert not at.exception

    labels = {m.label for m in at.metric}
    assert {"Revenue", "Operating income", "Headcount"} <= labels
    assert len(at.dataframe) >= 4

    # The deck is gated on Gate 1, so approve before it can be built.
    reviewer = next(t for t in at.text_input if t.label == "Your name")
    at = reviewer.set_value("Test Reviewer").run()
    at = next(b for b in at.button if b.label == "Sign off").click().run()

    button = next(b for b in at.button if b.label == "Build the deck")
    at = button.click().run()
    assert not at.exception

    downloads = at.get("download_button")
    assert downloads and downloads[0].label.startswith("variance-review-")
    # The figure count sits beside each download in the merged Deliverables
    # section, rendered as markdown rather than a caption.
    text = " ".join(m.value for m in at.markdown)
    assert "traced figures" in text


# --------------------------------------------------------------------------
# 5. portability: text I/O must not depend on the machine's locale
# --------------------------------------------------------------------------
def test_no_text_file_is_opened_without_an_explicit_encoding():
    """Regression for a defect that only appears on Windows.

    `open(path)` in text mode uses the LOCALE encoding: UTF-8 on Linux, cp1252
    on a default Windows install. Reading a source file containing an em dash
    or an arrow therefore works in CI and raises UnicodeDecodeError on the
    developer's machine.

    The worse case is silent. `run_agent.py` contains "·" (UTF-8 C2 B7), which
    cp1252 decodes without complaint as "Â·" -- so the model-boundary scan in
    test_agent_phase2 was reading mojibake and passing. Not a crash: a wrong
    read that still returns a string, on a test whose entire job is to inspect
    source text.

    Same shape as the build-hash defect: correct on the machine it was written
    on, wrong everywhere else, and invisible until someone ran it elsewhere.
    """
    import ast

    # Parsed, not grepped: a regex matches the word "open(" inside this very
    # docstring, which is how the first version of this test failed itself.
    roots = [os.path.join(SRC, "agent"), os.path.join(HERE, "..", "app"), HERE]
    offenders = []

    for root in roots:
        for dirpath, _dirs, files in os.walk(root):
            if "__pycache__" in dirpath:
                continue
            for fname in sorted(f for f in files if f.endswith(".py")):
                path = os.path.join(dirpath, fname)
                tree = ast.parse(open(path, encoding="utf-8").read())
                for node in ast.walk(tree):
                    if not (isinstance(node, ast.Call)
                            and isinstance(node.func, ast.Name)
                            and node.func.id == "open"):
                        continue
                    if any(k.arg == "encoding" for k in node.keywords):
                        continue
                    mode = ""
                    if len(node.args) > 1 and isinstance(node.args[1], ast.Constant):
                        mode = str(node.args[1].value)
                    for k in node.keywords:
                        if k.arg == "mode" and isinstance(k.value, ast.Constant):
                            mode = str(k.value.value)
                    if "b" in mode:                     # binary needs no encoding
                        continue
                    offenders.append(f"{fname}:{node.lineno}")

    assert not offenders, (
        "text-mode open() without an explicit encoding - these read as cp1252 "
        "on Windows: " + ", ".join(offenders[:12])
    )


def test_agent_sources_are_valid_utf8_and_decodable():
    """Belt and braces: the files themselves must be UTF-8."""
    agent_dir = os.path.join(SRC, "agent")
    for fname in sorted(os.listdir(agent_dir)):
        if fname.endswith(".py"):
            raw = open(os.path.join(agent_dir, fname), "rb").read()
            raw.decode("utf-8")     # raises if not valid UTF-8


# --------------------------------------------------------------------------
# 6. the agent page follows the sidebar month
# --------------------------------------------------------------------------
def _agent_app(monkeypatch, timeout=240):
    pytest.importorskip("streamlit")
    from streamlit.testing.v1 import AppTest

    for env in ("OPENAI_API_KEY", "ANTHROPIC_API_KEY"):
        monkeypatch.delenv(env, raising=False)
    sys.path.insert(0, os.path.join(HERE, "..", "eval"))
    at = AppTest.from_file(os.path.join(HERE, "..", "app", "app.py"),
                           default_timeout=timeout)
    at.run()
    return at.sidebar.radio[0].set_value("Close-cycle agent").run()


def test_the_agent_page_analyses_the_month_selected_in_the_sidebar(monkeypatch):
    """Every other page followed the sidebar month; the agent page did not.

    It was pinned to whichever runs happened to be cached, so selecting a
    different month changed the whole app except the agent -- which reads as
    the agent being disconnected from the tool rather than part of it. A
    deterministic run costs about 40ms and no credential, so there was never a
    reason for it not to follow.
    """
    at = _agent_app(monkeypatch)

    seen = {}
    for month in ("2025-09-01", "2024-06-01", "2025-12-01"):
        at = at.sidebar.selectbox[0].set_value(month).run()
        assert not at.exception, f"{month}: {at.exception}"

        revenue = next(m.value for m in at.metric if m.label == "Revenue")
        seen[month] = revenue

        header = [m.value for m in at.markdown if "Analyzing" in m.value]
        assert header, f"{month}: the page must name the month it is analysing"

    assert len(set(seen.values())) == 3, (
        f"the agent must produce different figures per month, got {seen}"
    )


def test_the_deck_produced_from_the_page_is_for_the_selected_month(monkeypatch):
    """The whole chain -- package, narrative, Gate 1, deck -- follows one
    month, so an approval can never attach to a different period's deck."""
    at = _agent_app(monkeypatch)
    at = at.sidebar.selectbox[0].set_value("2024-06-01").run()

    reviewer = next(t for t in at.text_input if t.label == "Your name")
    at = reviewer.set_value("Test Reviewer").run()
    at = next(b for b in at.button if b.label == "Sign off").click().run()
    assert any("Signed off" in s.value for s in at.success)

    build = next(b for b in at.button if b.label == "Build the deck")
    at = build.click().run()
    assert not at.exception

    downloads = at.get("download_button")
    assert downloads and "2024-06" in downloads[0].label, (
        f"the deck must be for the selected month, got {downloads[0].label}"
    )


def test_the_deck_is_withheld_until_it_is_signed_off(monkeypatch):
    """An unapproved deck is precisely the artifact that should not leave the
    building, so the download must not exist before approval."""
    at = _agent_app(monkeypatch)
    at = at.sidebar.selectbox[0].set_value("2025-03-01").run()

    assert not any(b.label == "Build the deck" for b in at.button)
    assert any("signed off above" in m.value for m in at.markdown)


def test_the_mode_selector_offers_exactly_two_modes(monkeypatch):
    """A third "verified replay" tab rendered the SAME package as the standard
    close with an integrity banner most readers could not interpret. Three
    modes where two looked identical cost more in confusion than the point was
    worth, so the tab was removed and the replay machinery kept -- it is still
    covered by the round-trip and tamper tests in phase 5."""
    at = _agent_app(monkeypatch)
    mode = next(r for r in at.radio
                if r.options and "Standard monthly close" in r.options)
    assert list(mode.options) == ["Standard monthly close", "Ask a question"]

    at = mode.set_value("Ask a question").run()
    assert not at.exception
    assert not any(sb.label == "Saved run" for sb in at.selectbox)


# --------------------------------------------------------------------------
# 6. a hex digest is not a figure
# --------------------------------------------------------------------------
HEX_IN_TEXT = [
    "Approved for publication by Ray · artifact d9434ff7908b617d",
    "run_id 3a8c3e208c6d",
    "artifact a7908b",
    "build_hash 9825377cafc18903",
    "FY2025 plan",
]


@pytest.mark.parametrize("text", HEX_IN_TEXT)
def test_a_hex_digest_is_not_read_as_a_figure(text):
    """Regression found by verifying three real decks against the data.

    The approval hash on one deck was `d9434ff7908b617d`, which contains
    "7908b". The trailing-boundary rule passed it -- the "b" is followed by a
    digit, not a letter -- so it parsed as $7.908 TRILLION and appeared as a
    figure with no computed source.

    The worst property of this bug is that it is data-dependent: it fires only
    when a random digest happens to contain a digit run followed by k/m/b. Two
    of the three decks were clean; the third was not. A leading boundary means
    a figure must START at one, so a digest yields nothing.
    """
    from guardrails.numeric_audit import extract_mentions

    assert extract_mentions(text) == [], f"{text!r} produced a phantom figure"


def test_the_appendix_hash_never_becomes_a_phantom_figure(run, goal, candidate):
    """End-to-end: audit the appendix slide of a deck whose approval hash is
    forced to the value that triggered the original defect."""
    from agent.deck import build_deck
    from agent.facts import fact_pack_from_ledger
    from agent.gates import PublicationPacket
    from guardrails import numeric_audit as na

    # The hash is patched at its SOURCE rather than overwritten on the
    # approval: overwriting it afterwards makes the approval refer to a
    # different artifact, which correctly voids it -- the binding rule doing
    # its job, and not what this test is about.
    import agent.gates as gates_mod

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(gates_mod, "artifact_hash",
                        lambda *_a, **_k: "d9434ff7908b617d")
    try:
        packet = PublicationPacket(run, candidate)
        packet.approve("Ray", "verification")
        assert packet.published, "the approval must be current for this test"
        d = build_deck(run, goal, candidate, packet)
    finally:
        monkeypatch.undo()
    pack = fact_pack_from_ledger(run, goal)

    appendix = d.prs.slides[len(d.prs.slides._sldIdLst) - 1]
    text = "\n".join(sh.text_frame.text for sh in appendix.shapes
                     if sh.has_text_frame)
    assert "d9434ff7908b617d" in text

    result = na.audit(text, pack)
    phantom = [v for v in result.violations if abs(v.value) > 1e12]
    assert not phantom, f"the hash was read as a figure: {phantom}"
